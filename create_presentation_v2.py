#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Presentation V2 - "Tabiiy resurslar va iqtisodiy rivojlanish"
A cohesive, professional editorial deck (navy + gold theme).

Design system
-------------
- One consistent palette across all 11 slides (deep navy ink + warm gold accent).
- Full-bleed backgrounds composited with PIL: smooth diagonal gradient,
  optional feathered photographic panel with a navy tonal overlay, soft gold glow.
- Typography: Georgia (titles, serif/editorial) + Arial (labels & body, sans).
- Each slide carries a gold kicker label, a thin rule, a footer and a page number.
- Restrained vector accents only (thin rules, small markers) - no loud boxes.
"""

import os
from PIL import Image, ImageDraw, ImageFilter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "assets", "img")
BG = os.path.join(BASE, "assets", "bg")
os.makedirs(BG, exist_ok=True)

# --------------------------------------------------------------------------
# Palette
# --------------------------------------------------------------------------
INK        = (10, 18, 32)       # deepest background
NAVY       = (17, 31, 52)       # panel navy
NAVY_LIT   = (28, 49, 79)       # lighter navy for depth
GOLD       = (201, 162, 75)     # primary accent
GOLD_LIT   = (227, 199, 122)    # bright gold
PAPER      = (244, 241, 234)    # primary text
MUTE       = (162, 178, 196)    # secondary text
LINE       = (60, 84, 116)      # hairline

# px canvas (16:9)
W, H = 1920, 1080

# slide (EMU) dimensions
SW = Inches(13.333)
SH = Inches(7.5)

# px -> inches mapping for placing elements that align with the background
def px_in(px):
    return Inches(px / W * 13.333)


# ==========================================================================
# PIL background compositing
# ==========================================================================
def _lerp(a, b, t):
    return tuple(int(a[i] + (b[i] - a[i]) * t) for i in range(3))


def diagonal_gradient(size, c1, c2):
    """Smooth diagonal gradient from c1 (top-left) to c2 (bottom-right)."""
    w, h = size
    base = Image.new("RGB", size, c1)
    # build a small gradient then resize for speed/smoothness
    grad = Image.new("L", (w, 1))
    for x in range(w):
        grad.putpixel((x, 0), int(255 * x / max(1, w - 1)))
    grad = grad.resize((w, h))
    # diagonal: average horizontal and vertical ramps
    gv = Image.new("L", (1, h))
    for y in range(h):
        gv.putpixel((0, y), int(255 * y / max(1, h - 1)))
    gv = gv.resize((w, h))
    diag = Image.blend(grad, gv, 0.5)
    top = Image.new("RGB", size, c2)
    return Image.composite(top, base, diag)


def radial_glow(size, center, radius, color, max_alpha):
    """A soft radial glow layer (RGBA)."""
    w, h = size
    glow = Image.new("L", size, 0)
    d = ImageDraw.Draw(glow)
    cx, cy = center
    d.ellipse([cx - radius, cy - radius, cx + radius, cy + radius], fill=max_alpha)
    glow = glow.filter(ImageFilter.GaussianBlur(radius * 0.45))
    layer = Image.new("RGBA", size, color + (0,))
    layer.putalpha(glow)
    return layer


def cover_crop(img, target_w, target_h):
    """Resize+crop image to exactly cover target box (like CSS background-size:cover)."""
    iw, ih = img.size
    scale = max(target_w / iw, target_h / ih)
    nw, nh = int(iw * scale + 0.5), int(ih * scale + 0.5)
    img = img.resize((nw, nh), Image.LANCZOS)
    left = (nw - target_w) // 2
    top = (nh - target_h) // 3  # bias slightly toward upper third
    return img.crop((left, top, left + target_w, top + target_h))


def make_background(name, photo=None, layout="panel", overlay=0.45, glow_corner="tr"):
    """
    Compose a full slide background.
      layout = "panel"  -> photo on right ~46%, feathered into navy on the left
      layout = "hero"   -> full-bleed photo with strong navy gradient overlay
      layout = "plain"  -> gradient only (no photo)
    """
    bg = diagonal_gradient((W, H), INK, NAVY_LIT)

    # subtle gold glow for richness
    if glow_corner == "tr":
        center = (int(W * 0.86), int(H * 0.18))
    elif glow_corner == "bl":
        center = (int(W * 0.14), int(H * 0.82))
    else:
        center = (int(W * 0.5), int(H * 0.5))
    bg = bg.convert("RGBA")
    bg.alpha_composite(radial_glow((W, H), center, int(W * 0.34), GOLD, 46))

    if photo and layout in ("panel", "hero"):
        src = Image.open(os.path.join(IMG, photo)).convert("RGB")

        if layout == "hero":
            pic = cover_crop(src, W, H)
            # tonal cohesion: blend toward navy
            pic = Image.blend(pic, Image.new("RGB", (W, H), NAVY), overlay)
            pic = pic.convert("RGBA")
            # darken gradient: darker at bottom-left where text lives
            shade = Image.new("L", (W, H), 0)
            ds = ImageDraw.Draw(shade)
            for y in range(H):
                # stronger at bottom
                a = int(150 * (y / H) ** 1.4)
                ds.line([(0, y), (W, y)], fill=a)
            dark = Image.new("RGBA", (W, H), INK + (0,))
            dark.putalpha(shade)
            pic.alpha_composite(dark)
            # left-side vignette for title readability
            lshade = Image.new("L", (W, 1), 0)
            for x in range(W):
                a = int(170 * max(0.0, 1 - x / (W * 0.62)))
                lshade.putpixel((x, 0), a)
            lshade = lshade.resize((W, H))
            lv = Image.new("RGBA", (W, H), INK + (0,))
            lv.putalpha(lshade)
            pic.alpha_composite(lv)
            bg.alpha_composite(pic)

        else:  # panel
            pw = int(W * 0.46)
            ph = H
            px = W - pw
            pic = cover_crop(src, pw, ph)
            pic = Image.blend(pic, Image.new("RGB", (pw, ph), NAVY), overlay)
            pic = pic.convert("RGBA")
            # feather the left edge of the panel into the navy
            mask = Image.new("L", (pw, ph), 255)
            feather = int(pw * 0.34)
            mcol = Image.new("L", (pw, 1), 255)
            for x in range(pw):
                if x < feather:
                    mcol.putpixel((x, 0), int(255 * (x / feather) ** 1.3))
            mcol = mcol.resize((pw, ph))
            pic.putalpha(mcol)
            # bottom darken for footer legibility
            sh = Image.new("L", (1, ph), 0)
            for y in range(ph):
                sh.putpixel((0, y), int(120 * (y / ph) ** 2))
            sh = sh.resize((pw, ph))
            dk = Image.new("RGBA", (pw, ph), INK + (0,))
            dk.putalpha(sh)
            pic.alpha_composite(dk)
            bg.alpha_composite(pic, (px, 0))
            # thin gold seam between text area and photo
            seam = Image.new("RGBA", (W, H), (0, 0, 0, 0))
            sd = ImageDraw.Draw(seam)
            sx = px + int(feather * 0.62)
            sd.line([(sx, int(H * 0.16)), (sx, int(H * 0.84))], fill=GOLD + (120,), width=2)
            bg.alpha_composite(seam)

    out = os.path.join(BG, f"{name}.jpg")
    bg.convert("RGB").save(out, "JPEG", quality=86, optimize=True, progressive=True)
    return out


# ==========================================================================
# python-pptx helpers
# ==========================================================================
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_bg(slide, png_path):
    slide.shapes.add_picture(png_path, 0, 0, width=SW, height=SH)


def _set_run(r, text, size, color, bold=False, italic=False, font="Arial", spacing=None):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = RGBColor(*color)
    if spacing is not None:
        # letter spacing in points (1/100 pt units -> use 'spc' on rPr in 1/100 pt)
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def text_block(slide, left, top, width, height, paragraphs, anchor=MSO_ANCHOR.TOP):
    """
    paragraphs: list of dicts:
      {runs:[(text,size,color,bold,italic,font,spacing)], align, space_after, space_before, line}
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
        if "space_before" in para:
            p.space_before = Pt(para["space_before"])
        if "line" in para:
            p.line_spacing = para["line"]
        for rspec in para["runs"]:
            r = p.add_run()
            _set_run(r, *rspec)
    return tb


def rect(slide, left, top, width, height, color, alpha=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor(*color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    if alpha is not None:
        _set_alpha(sp, alpha)
    return sp


def _set_alpha(shape, alpha_pct):
    """Set fill transparency (alpha_pct = 0..100 opacity)."""
    sp = shape.fill._xPr.find(qn("a:solidFill"))
    srgb = sp.find(qn("a:srgbClr"))
    a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
    srgb.append(a)


# kicker label + title header, consistent on every slide
def header(slide, kicker, title_lines, title_size=40, top=Inches(0.72), color=PAPER):
    # gold kicker
    text_block(slide, Inches(0.9), top, Inches(8.5), Inches(0.4), [
        {"runs": [(kicker, 14, GOLD_LIT, True, False, "Arial", 3.0)], "align": PP_ALIGN.LEFT},
    ])
    # thin rule
    rect(slide, Inches(0.92), top + Inches(0.42), Inches(0.62), Pt(2.4), GOLD)
    # title (Georgia serif)
    runs_paras = []
    for j, line in enumerate(title_lines):
        runs_paras.append({
            "runs": [(line, title_size, color, True, False, "Georgia")],
            "align": PP_ALIGN.LEFT, "line": 1.04,
            "space_after": 0,
        })
    text_block(slide, Inches(0.88), top + Inches(0.6), Inches(11), Inches(1.8), runs_paras)


def footer(slide, page):
    text_block(slide, Inches(0.9), Inches(7.02), Inches(8), Inches(0.35), [
        {"runs": [("TABIIY RESURSLAR VA IQTISODIY RIVOJLANISH", 9, MUTE, False, False, "Arial", 2.0)]},
    ])
    text_block(slide, Inches(11.3), Inches(7.02), Inches(1.2), Inches(0.35), [
        {"runs": [(f"{page:02d}", 11, GOLD_LIT, True, False, "Arial"),
                  ("  /  11", 11, MUTE, False, False, "Arial")], "align": PP_ALIGN.RIGHT},
    ])
    rect(slide, Inches(0.9), Inches(6.95), Inches(11.53), Pt(0.75), LINE)


def new_slide(bg_png):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, bg_png)
    return s


# small gold diamond/marker for list items
def marker(slide, left, top):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, Pt(8), Pt(8))
    d.fill.solid()
    d.fill.fore_color.rgb = RGBColor(*GOLD)
    d.line.fill.background()
    d.shadow.inherit = False


# numbered editorial row: big light gold numeral + heading + supporting line
def num_row(slide, left, top, width, n, heading, sub=None):
    text_block(slide, left, top - Inches(0.06), Inches(0.9), Inches(0.9), [
        {"runs": [(f"{n:02d}", 30, GOLD, True, False, "Georgia")]},
    ])
    text_block(slide, left + Inches(1.0), top, width - Inches(1.0), Inches(0.9),
               ([{"runs": [(heading, 20, PAPER, True, False, "Arial")], "space_after": 2}] +
                ([{"runs": [(sub, 13, MUTE, False, False, "Arial")], "line": 1.1}] if sub else [])))
    rect(slide, left + Inches(1.0), top + Inches(0.82), Inches(4.6), Pt(0.75), LINE)


# ==========================================================================
# Generate backgrounds
# ==========================================================================
print("Composing backgrounds...")
bg_s1  = make_background("s01", "gold.jpg",       layout="hero", overlay=0.50, glow_corner="tr")
bg_s2  = make_background("s02", "mountains.jpg",  layout="panel", overlay=0.42, glow_corner="bl")
bg_s3  = make_background("s03", "industry.jpg",   layout="panel", overlay=0.46, glow_corner="tr")
bg_s4  = make_background("s04", "coins.jpg",      layout="panel", overlay=0.52, glow_corner="bl")
bg_s5  = make_background("s05", None,             layout="plain", glow_corner="tr")
bg_s6  = make_background("s06", "fjord.jpg",      layout="panel", overlay=0.40, glow_corner="tr")
bg_s7  = make_background("s07", "industry.jpg",   layout="panel", overlay=0.55, glow_corner="bl")
bg_s8  = make_background("s08", "samarkand.jpg",  layout="panel", overlay=0.40, glow_corner="tr")
bg_s9  = make_background("s09", "circuit.jpg",    layout="panel", overlay=0.48, glow_corner="bl")
bg_s10 = make_background("s10", "meeting.jpg",    layout="panel", overlay=0.48, glow_corner="tr")
bg_s11 = make_background("s11", "city_night.jpg", layout="hero", overlay=0.48, glow_corner="tr")

# content text width when a photo panel sits on the right
LX = Inches(0.9)         # left margin
TXTW = Inches(6.7)       # text width over navy area (panel slides)
FULLW = Inches(11.5)

# ==========================================================================
# SLIDE 1 - Hook
# ==========================================================================
s = new_slide(bg_s1)
text_block(s, LX, Inches(1.7), Inches(9), Inches(0.5), [
    {"runs": [("KIRISH  \u2014  MUNOZARA SAVOLI", 15, GOLD_LIT, True, False, "Arial", 3.0)]},
])
rect(s, Inches(0.92), Inches(2.18), Inches(0.7), Pt(3), GOLD)
text_block(s, LX, Inches(2.45), Inches(11.3), Inches(3), [
    {"runs": [("Agar ertaga O\u2018zbekistonda dunyodagi eng katta",
               38, PAPER, True, False, "Georgia")], "line": 1.08},
    {"runs": [("oltin yoki neft koni topilsa \u2014 biz avtomatik",
               38, PAPER, True, False, "Georgia")], "line": 1.08},
    {"runs": [("ravishda ", 38, PAPER, True, False, "Georgia"),
              ("boy davlatga", 38, GOLD_LIT, True, False, "Georgia"),
              (" aylanamizmi?", 38, PAPER, True, False, "Georgia")], "line": 1.08},
])
text_block(s, LX, Inches(5.7), Inches(9), Inches(0.6), [
    {"runs": [("Tinglovchilardan javoblarini so\u2018rang \u2014 munozarani shu yerdan boshlaymiz.",
               16, MUTE, False, True, "Arial")]},
])
footer(s, 1)

# ==========================================================================
# SLIDE 2 - Tabiiy resurslar nima?
# ==========================================================================
s = new_slide(bg_s2)
header(s, "ASOSIY TUSHUNCHA", ["Tabiiy resurslar nima?"])
items = [
    "Neft", "Tabiiy gaz", "Oltin", "Mis",
    "Ko\u2018mir", "Suv", "Qishloq xo\u2018jaligi yerlari",
]
y = 2.55
col_x = [LX, Inches(3.7)]
for i, it in enumerate(items):
    cx = col_x[i % 2]
    row = i // 2
    yy = Inches(y + row * 0.72)
    marker(s, cx + Inches(0.02), yy + Inches(0.07))
    text_block(s, cx + Inches(0.32), yy, Inches(3.1), Inches(0.6), [
        {"runs": [(it, 19, PAPER, False, False, "Arial")]},
    ])
text_block(s, LX, Inches(6.05), TXTW, Inches(0.7), [
    {"runs": [("Tabiat bizga bergan boylik \u2014 ulardan qanday foydalanishimiz esa o\u2018zimizga bog\u2018liq.",
               14, MUTE, False, True, "Arial")], "line": 1.2},
])
footer(s, 2)

# ==========================================================================
# SLIDE 3 - Afzalliklar
# ==========================================================================
s = new_slide(bg_s3)
header(s, "IMKONIYATLAR", ["Resurslarning afzalliklari"])
adv = [
    ("Ish o\u2018rinlari", "Sanoat va xizmat sohalarida bandlik yaratadi"),
    ("Eksport daromadi", "Tashqi savdo orqali valyuta oqimi"),
    ("Davlat budjeti", "Soliq va to\u2018lovlar bilan to\u2018ldiriladi"),
    ("Infratuzilma", "Yo\u2018l, energiya va shaharlar rivojlanadi"),
]
y = 2.5
for i, (h, sub) in enumerate(adv):
    yy = Inches(y + i * 1.02)
    marker(s, LX + Inches(0.02), yy + Inches(0.12))
    text_block(s, LX + Inches(0.34), yy, TXTW - Inches(0.34), Inches(0.95), [
        {"runs": [(h, 19, PAPER, True, False, "Arial")], "space_after": 2},
        {"runs": [(sub, 13.5, MUTE, False, False, "Arial")], "line": 1.12},
    ])
text_block(s, LX, Inches(6.5), Inches(7), Inches(0.4), [
    {"runs": [("MISOLLAR:   ", 12, GOLD_LIT, True, False, "Arial", 2.0),
              ("Saudi Arabia   \u00B7   Norvegiya", 13, PAPER, False, False, "Arial")]},
])
footer(s, 3)

# ==========================================================================
# SLIDE 4 - Resource Curse
# ==========================================================================
s = new_slide(bg_s4)
header(s, "MUAMMO", ["Lekin muammo ham bor\u2026"])
text_block(s, LX, Inches(2.55), TXTW, Inches(1.4), [
    {"runs": [("\u201CResource Curse\u201D", 46, GOLD_LIT, True, False, "Georgia")], "space_after": 2},
    {"runs": [("Resurslar la\u2019nati \u2014 iqtisodiyotdagi mashhur tushuncha",
               16, MUTE, False, True, "Arial")]},
])
text_block(s, LX, Inches(4.35), TXTW, Inches(1.6), [
    {"runs": [("Ba\u2019zi davlatlar tabiiy resurslarga juda boy,",
               19, PAPER, False, False, "Arial")], "space_after": 6, "line": 1.15},
    {"runs": [("ammo aholisi kambag\u2018alligicha qoladi.",
               19, PAPER, False, False, "Arial")], "line": 1.15},
])
text_block(s, LX, Inches(5.95), Inches(5), Inches(0.7), [
    {"runs": [("Nega shunday bo\u2018ladi?", 22, GOLD, True, False, "Georgia")]},
])
footer(s, 4)

# ==========================================================================
# SLIDE 5 - Sabablar (plain, no photo -> use full width, two columns)
# ==========================================================================
s = new_slide(bg_s5)
header(s, "TAHLIL", ["Resurslar la\u2019natining sabablari"])
causes = [
    ("Korrupsiya", "Daromadlar tor doiraga yo\u2018naltiriladi"),
    ("Noto\u2018g\u2018ri boshqaruv", "Strategiya va shaffoflik yetishmaydi"),
    ("Bitta sohaga tayanish", "Iqtisodiyot diversifikatsiya qilinmaydi"),
    ("Sust innovatsiya", "Texnologiya va ta\u2019limga investitsiya kam"),
]
positions = [(LX, 2.7), (Inches(6.9), 2.7), (LX, 4.7), (Inches(6.9), 4.7)]
for i, (h, sub) in enumerate(causes):
    px_, py_ = positions[i]
    num_row(s, px_, Inches(py_), Inches(5.6), i + 1, h, sub)
footer(s, 5)

# ==========================================================================
# SLIDE 6 - Norway
# ==========================================================================
s = new_slide(bg_s6)
header(s, "MUVAFFAQIYAT MISOLI", ["Norvegiya"])
steps = [
    ("Neft topildi", "1969-yil, Shimoliy dengizda"),
    ("Maxsus fond", "Daromadlar suveren fondga yig\u2018ildi"),
    ("Investitsiya", "Ta\u2019lim va infratuzilmaga yo\u2018naltirildi"),
]
y = 2.55
for i, (h, sub) in enumerate(steps):
    num_row(s, LX, Inches(y + i * 1.12), TXTW, i + 1, h, sub)
text_block(s, LX, Inches(6.05), TXTW, Inches(0.8), [
    {"runs": [("NATIJA:   ", 12, GOLD_LIT, True, False, "Arial", 2.0),
              ("dunyodagi eng yuqori yashash standartlaridan biri.",
               15, PAPER, False, False, "Arial")], "line": 1.2},
])
footer(s, 6)

# ==========================================================================
# SLIDE 7 - Venezuela
# ==========================================================================
s = new_slide(bg_s7)
header(s, "OGOHLANTIRUVCHI MISOL", ["Venesuela"])
text_block(s, LX, Inches(2.5), TXTW, Inches(0.9), [
    {"runs": [("Juda katta neft zaxiralari \u2014 ammo chuqur iqtisodiy inqirozlar.",
               17, PAPER, False, False, "Arial")], "line": 1.2},
])
vcz = [
    ("Noto\u2018g\u2018ri siyosat", "Qisqa muddatli, nomutanosib qarorlar"),
    ("Yuqori inflyatsiya", "Pul qadrsizlanishi, narxlar o\u2018sishi"),
    ("Neftga bog\u2018liqlik", "Boshqa sohalar rivojlanmay qoldi"),
]
y = 3.45
for i, (h, sub) in enumerate(vcz):
    yy = Inches(y + i * 1.02)
    marker(s, LX + Inches(0.02), yy + Inches(0.12))
    text_block(s, LX + Inches(0.34), yy, TXTW - Inches(0.34), Inches(0.95), [
        {"runs": [(h, 19, PAPER, True, False, "Arial")], "space_after": 2},
        {"runs": [(sub, 13.5, MUTE, False, False, "Arial")], "line": 1.12},
    ])
footer(s, 7)

# ==========================================================================
# SLIDE 8 - O'zbekiston
# ==========================================================================
s = new_slide(bg_s8)
header(s, "BIZNING DAVLAT", ["O\u2018zbekiston misolida"])
res_uz = ["Oltin", "Mis", "Uran", "Tabiiy gaz"]
y = 2.6
for i, it in enumerate(res_uz):
    cx = col_x[i % 2]
    row = i // 2
    yy = Inches(y + row * 0.78)
    marker(s, cx + Inches(0.02), yy + Inches(0.08))
    text_block(s, cx + Inches(0.32), yy, Inches(3.0), Inches(0.6), [
        {"runs": [(it, 20, PAPER, False, False, "Arial")]},
    ])
text_block(s, LX, Inches(4.7), TXTW, Inches(1.2), [
    {"runs": [("Bu resurslardan qanday qilib", 22, PAPER, True, False, "Georgia")], "line": 1.1, "space_after": 2},
    {"runs": [("samarali", 22, GOLD_LIT, True, False, "Georgia"),
              (" foydalanish mumkin?", 22, PAPER, True, False, "Georgia")], "line": 1.1},
])
footer(s, 8)

# ==========================================================================
# SLIDE 9 - Kelajak
# ==========================================================================
s = new_slide(bg_s9)
header(s, "KELAJAK", ["Eng qimmat resurs nima bo\u2018ladi?"])
fut = [
    ("Bilim", "Ta\u2019lim va malaka \u2014 yangi boylik manbai"),
    ("Texnologiya", "Innovatsiya iqtisodiy o\u2018sishni belgilaydi"),
    ("Inson kapitali", "Odamlar \u2014 davlatning asosiy resursi"),
]
y = 2.5
for i, (h, sub) in enumerate(fut):
    num_row(s, LX, Inches(y + i * 1.02), TXTW, i + 1, h, sub)
text_block(s, LX, Inches(5.85), TXTW, Inches(1.0), [
    {"runs": [("\u201C21-asrning eng muhim resursi \u2014 insonning bilimi.\u201D",
               16, GOLD_LIT, False, True, "Georgia")], "line": 1.2},
])
footer(s, 9)

# ==========================================================================
# SLIDE 10 - Muhokama
# ==========================================================================
s = new_slide(bg_s10)
header(s, "MUNOZARA", ["Muhokama uchun savollar"])
qs = [
    "Tabiiy resurslarsiz boy bo\u2018lish mumkinmi?",
    "O\u2018zbekistonda qaysi resursdan yetarlicha foydalanilmayapti?",
    "Bosh vazir bo\u2018lsangiz, resurs daromadlarini nimaga sarflardingiz?",
]
y = 2.7
for i, q in enumerate(qs):
    yy = Inches(y + i * 1.18)
    text_block(s, LX, yy - Inches(0.05), Inches(0.9), Inches(0.9), [
        {"runs": [(f"{i+1:02d}", 26, GOLD, True, False, "Georgia")]},
    ])
    text_block(s, LX + Inches(0.95), yy, TXTW - Inches(0.95), Inches(1.0), [
        {"runs": [(q, 18, PAPER, False, False, "Arial")], "line": 1.18},
    ])
    rect(s, LX + Inches(0.95), yy + Inches(0.86), Inches(4.4), Pt(0.75), LINE)
footer(s, 10)

# ==========================================================================
# SLIDE 11 - Xulosa
# ==========================================================================
s = new_slide(bg_s11)
text_block(s, LX, Inches(1.55), Inches(9), Inches(0.5), [
    {"runs": [("XULOSA", 15, GOLD_LIT, True, False, "Arial", 3.5)]},
])
rect(s, Inches(0.92), Inches(2.02), Inches(0.7), Pt(3), GOLD)
text_block(s, LX, Inches(2.3), Inches(11.4), Inches(2.2), [
    {"runs": [("Resurslar davlatni ", 33, PAPER, True, False, "Georgia"),
              ("boy", 33, GOLD_LIT, True, False, "Georgia"),
              (" qilishi mumkin,", 33, PAPER, True, False, "Georgia")], "line": 1.1, "space_after": 2},
    {"runs": [("lekin uni uzoq muddatda ", 33, PAPER, True, False, "Georgia"),
              ("farovon", 33, GOLD_LIT, True, False, "Georgia"),
              (" qiladigan narsa \u2014", 33, PAPER, True, False, "Georgia")], "line": 1.1, "space_after": 2},
    {"runs": [("yaxshi institutlar, ta\u2019lim va innovatsiya.",
               33, PAPER, True, False, "Georgia")], "line": 1.1},
])
text_block(s, LX, Inches(5.15), Inches(11), Inches(0.9), [
    {"runs": [("\u201CResources can make a country rich, but only good institutions, "
               "education, and innovation can make it prosperous in the long run.\u201D",
               14, MUTE, False, True, "Georgia")], "line": 1.25},
])
text_block(s, LX, Inches(6.2), Inches(11), Inches(0.5), [
    {"runs": [("Bu mavzu iqtisodiyot, siyosat, tarix va kelajak texnologiyalarini bir nuqtada birlashtiradi.",
               13, MUTE, False, False, "Arial")]},
])
footer(s, 11)

# ==========================================================================
# Save
# ==========================================================================
out = os.path.join(BASE, "presentation_v2.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
